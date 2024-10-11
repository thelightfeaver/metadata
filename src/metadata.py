from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pprint import pprint


class Metadata:
    """
    The Metadata class provides methods to read and write metadata for various file types including docx, pptx, and xlsx.
    Attributes:
        _filename (str): The name of the file to read/write metadata.
        _structures (list): A list of metadata attributes available in the file.
    Methods:
        __init__(filename):
            Initializes the Metadata object with the given filename and retrieves the metadata structures.
        read_docx_metadata():
            Reads metadata from a docx file and returns it as a dictionary.
        write_docx_metadata(data):
            Writes metadata to a docx file using the provided dictionary.
        _get_data():
            Retrieves the metadata structures based on the file extension and populates the _structures attribute.
        read_pptx_metadata():
            Reads metadata from a pptx file and returns it as a dictionary.
        write_pptx_metadata(data):
            Writes metadata to a pptx file using the provided dictionary.
        read_xlsx_metadata():
            Reads metadata from a xlsx file and returns it as a dictionary.
        write_xlsx_metadata(data):
            Writes metadata to a xlsx file using the provided dictionary.
    """

    def __init__(self, filename):
        self._filename = filename
        self._structures = list()
        self._get_data()

    def _get_data(self):
        """Get the structures of the metadata"""

        file_extension = self._filename.split(".")[-1]
        if file_extension in ["docx", "docm", "dotx", "dotm"]:
            file = Document(self._filename)
            for attr in dir(file.core_properties):
                if (
                    not callable(attr)
                    and not str(attr).startswith("__")
                    and not str(attr).startswith("_")
                ):
                    self._structures.append(attr)

        elif file_extension in ["pptx", "pptm", "potx", "potm", "ppsx", "ppsm"]:
            file = Presentation(self._filename)
            for attr in dir(file.core_properties):
                if (
                    not callable(attr)
                    and not str(attr).startswith("__")
                    and not str(attr).startswith("_")
                ):
                    self._structures.append(attr)
        elif file_extension in ["xlsx", "xlsm", "xltx", "xltm"]:
            file = load_workbook(self._filename)

            for attr in dir(file.properties):
                if (
                    not callable(attr)
                    and not str(attr).startswith("__")
                    and not str(attr).startswith("_")
                ):
                    self._structures.append(attr)
        else:
            raise Exception("File extension not supported")

    def read_docx_metadata(self):
        """Reads metadata from a docx file"""
        doc = Document(self._filename)
        mt = dict()
        for attr in self._structures:
            mt[attr] = getattr(doc.core_properties, attr)
        return mt

    def write_docx_metadata(self, data):
        """
        Writes metadata to a docx file.

        Args:
            data (dict): A dictionary containing metadata attributes and values to be written to the file.
                         Possible keys include:
                         - 'author'
                         - 'category'
                         - 'comments'
                         - 'content_status'
                         - 'created'
                         - 'identifier'
                         - 'keywords'
                         - 'language'
                         - 'last_modified_by'
                         - 'last_printed'
                         - 'modified'
                         - 'revision'
                         - 'subject'
                         - 'title'
                         - 'version'
        """
        doc = Document(self._filename)
        for attr in data:
            setattr(doc.core_properties, attr, data[attr])
        doc.save(self._filename)

    def read_pptx_metadata(self):
        """Reads metadata from a pptx file"""
        ppt = Presentation(self._filename)
        mt = dict()
        for attr in self._structures:
            mt[attr] = getattr(ppt.core_properties, attr)
        return mt

    def write_pptx_metadata(self, data):
        """
        Writes metadata to a pptx file.

        Args:
            data (dict): A dictionary containing metadata attributes and values to be written to the file.
                         Possible keys include:
                         - 'author'
                         - 'category'
                         - 'comments'
                         - 'content_status'
                         - 'created'
                         - 'identifier'
                         - 'keywords'
                         - 'language'
                         - 'last_modified_by'
                         - 'last_printed'
                         - 'modified'
                         - 'revision'
                         - 'subject'
                         - 'title'
                         - 'version'
        """
        ppt = Presentation(self._filename)
        for attr in data:
            setattr(ppt.core_properties, attr, data[attr])
        ppt.save(self._filename)

    def read_xlsx_metadata(self):
        """Reads metadata from a xlsx file"""
        wb = load_workbook(self._filename)
        mt = dict()
        for attr in self._structures:
            mt[attr] = getattr(wb.properties, attr)
        return mt

    def write_xlsx_metadata(self, data):
        """
        Writes metadata to a xlsx file.

        Args:
            data (dict): A dictionary containing metadata attributes and values to be written to the file.
                         Possible keys include:
                         - 'category'
                         - 'contentStatus'
                         - 'created'
                         - 'creator'
                         - 'description'
                         - 'identifier'
                         - 'keywords'
                         - 'language'
                         - 'lastModifiedBy'
                         - 'lastPrinted'
                         - 'modified'
                         - 'revision'
                         - 'subject'
                         - 'title'
                         - 'version'
        """
        wb = load_workbook(self._filename)
        for attr in data:
            setattr(wb.properties, attr, data[attr])
        wb.save(self._filename)
